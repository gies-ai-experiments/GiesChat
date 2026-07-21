"""
Presentation management utilities for PowerPoint MCP Server.
Functions for creating, opening, saving, and managing presentations.
"""
from pptx import Presentation
from typing import Dict, List, Optional
import os

from gies_auth import current_user
from gies_sandbox import resolve as _sandbox


def create_presentation() -> Presentation:
    """
    Create a new PowerPoint presentation.
    
    Returns:
        A new Presentation object
    """
    return Presentation()


def open_presentation(file_path: str) -> Presentation:
    """
    Open an existing PowerPoint presentation.

    Args:
        file_path: Path to the PowerPoint file

    Returns:
        A Presentation object
    """
    file_path = _sandbox(file_path, current_user())
    return Presentation(file_path)


def create_presentation_from_template(template_path: str) -> Presentation:
    """
    Create a new PowerPoint presentation from a template file.
    
    Args:
        template_path: Path to the template .pptx file
        
    Returns:
        A new Presentation object based on the template
        
    Raises:
        FileNotFoundError: If the template file doesn't exist
        Exception: If the template file is corrupted or invalid
    """
    template_path = _sandbox(template_path, current_user(), allow_template=True)
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template file not found: {template_path}")
    
    if not template_path.lower().endswith(('.pptx', '.potx')):
        raise ValueError("Template file must be a .pptx or .potx file")
    
    try:
        # Load the template file as a presentation
        presentation = Presentation(template_path)
        return presentation
    except Exception as e:
        raise Exception(f"Failed to load template file '{template_path}': {str(e)}")


MIN_FONT_PT = 18


def _strip_empty_placeholders(presentation: Presentation) -> None:
    """Gies: delete placeholders the model never filled, so downloaded decks
    don't show "Click to add title"/"Click to add text" prompts overlapping
    the text boxes the model drew instead."""
    for slide in presentation.slides:
        for shape in list(slide.placeholders):
            if shape.has_text_frame and not shape.text_frame.text.strip():
                element = shape._element
                element.getparent().remove(element)


def _bump_small_fonts(presentation: Presentation) -> None:
    """Gies: enforce a floor on explicitly-set font sizes; models routinely
    pick 10-14pt body text that is unreadable on a projected slide."""
    from pptx.util import Pt

    floor = Pt(MIN_FONT_PT)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.font.size is not None and paragraph.font.size < floor:
                    paragraph.font.size = floor
                for run in paragraph.runs:
                    if run.font.size is not None and run.font.size < floor:
                        run.font.size = floor


def save_presentation(presentation: Presentation, file_path: str) -> str:
    """
    Save a PowerPoint presentation to a file.

    Args:
        presentation: The Presentation object
        file_path: Path where the file should be saved

    Returns:
        The file path where the presentation was saved
    """
    file_path = _sandbox(file_path, current_user())
    _strip_empty_placeholders(presentation)
    _bump_small_fonts(presentation)
    presentation.save(file_path)
    return file_path


def get_template_info(template_path: str) -> Dict:
    """
    Get information about a template file.
    
    Args:
        template_path: Path to the template .pptx file
        
    Returns:
        Dictionary containing template information
    """
    template_path = _sandbox(template_path, current_user(), allow_template=True)
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template file not found: {template_path}")
    
    try:
        presentation = Presentation(template_path)
        
        # Get slide layouts
        layouts = get_slide_layouts(presentation)
        
        # Get core properties
        core_props = get_core_properties(presentation)
        
        # Get slide count
        slide_count = len(presentation.slides)
        
        # Get file size
        file_size = os.path.getsize(template_path)
        
        return {
            "template_path": template_path,
            "file_size_bytes": file_size,
            "slide_count": slide_count,
            "layout_count": len(layouts),
            "slide_layouts": layouts,
            "core_properties": core_props
        }
    except Exception as e:
        raise Exception(f"Failed to read template info from '{template_path}': {str(e)}")


def get_presentation_info(presentation: Presentation) -> Dict:
    """
    Get information about a presentation.
    
    Args:
        presentation: The Presentation object
        
    Returns:
        Dictionary containing presentation information
    """
    try:
        # Get slide layouts
        layouts = get_slide_layouts(presentation)
        
        # Get core properties
        core_props = get_core_properties(presentation)
        
        # Get slide count
        slide_count = len(presentation.slides)
        
        return {
            "slide_count": slide_count,
            "layout_count": len(layouts),
            "slide_layouts": layouts,
            "core_properties": core_props,
            "slide_width": presentation.slide_width,
            "slide_height": presentation.slide_height
        }
    except Exception as e:
        raise Exception(f"Failed to get presentation info: {str(e)}")


def get_slide_layouts(presentation: Presentation) -> List[Dict]:
    """
    Get all available slide layouts in the presentation.
    
    Args:
        presentation: The Presentation object
        
    Returns:
        A list of dictionaries with layout information
    """
    layouts = []
    for i, layout in enumerate(presentation.slide_layouts):
        layout_info = {
            "index": i,
            "name": layout.name,
            "placeholder_count": len(layout.placeholders)
        }
        layouts.append(layout_info)
    return layouts


def set_core_properties(presentation: Presentation, title: str = None, subject: str = None,
                       author: str = None, keywords: str = None, comments: str = None) -> None:
    """
    Set core document properties.
    
    Args:
        presentation: The Presentation object
        title: Document title
        subject: Document subject
        author: Document author
        keywords: Document keywords
        comments: Document comments
    """
    core_props = presentation.core_properties
    
    if title is not None:
        core_props.title = title
    if subject is not None:
        core_props.subject = subject
    if author is not None:
        core_props.author = author
    if keywords is not None:
        core_props.keywords = keywords
    if comments is not None:
        core_props.comments = comments


def get_core_properties(presentation: Presentation) -> Dict:
    """
    Get core document properties.
    
    Args:
        presentation: The Presentation object
        
    Returns:
        Dictionary containing core properties
    """
    core_props = presentation.core_properties
    
    return {
        "title": core_props.title,
        "subject": core_props.subject,
        "author": core_props.author,
        "keywords": core_props.keywords,
        "comments": core_props.comments,
        "created": core_props.created.isoformat() if core_props.created else None,
        "last_modified_by": core_props.last_modified_by,
        "modified": core_props.modified.isoformat() if core_props.modified else None
    }