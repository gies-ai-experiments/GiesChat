import json
import os
import socket
import sys
import threading
import time

import httpx
import pytest
import uvicorn
from pptx import Presentation

_OWNED = {"gies_auth", "gies_sandbox", "gies_downloads", "gies_state",
          "gies_questions", "gies_uploads", "ppt_mcp_server", "gies_server",
          "utils", "tools"}


def _free_port():
    s = socket.socket()
    s.bind(("127.0.0.1", 0))
    port = s.getsockname()[1]
    s.close()
    return port


@pytest.fixture(scope="module")
def server(tmp_path_factory):
    tmp = tmp_path_factory.mktemp("pptx")
    tpl = tmp / "gies.pptx"
    Presentation().save(str(tpl))
    port = _free_port()
    os.environ.update(
        PPTX_MCP_KEY="testkey",
        PUBLIC_URL=f"http://127.0.0.1:{port}",
        GIES_TEMPLATE_PATH=str(tpl),
        PPT_TEMPLATE_PATH=str(tmp),
        PPTX_SANDBOX_ROOT=str(tmp / "decks"),
    )
    # Re-import the whole stack fresh so every module binds to the env above.
    for name in list(sys.modules):
        if name.split(".")[0] in _OWNED:
            del sys.modules[name]
    import gies_server

    config = uvicorn.Config(gies_server.asgi, host="127.0.0.1", port=port, log_level="error")
    srv = uvicorn.Server(config)
    thread = threading.Thread(target=srv.run, daemon=True)
    thread.start()
    for _ in range(100):
        try:
            socket.create_connection(("127.0.0.1", port), timeout=0.1).close()
            break
        except OSError:
            time.sleep(0.05)
    yield f"http://127.0.0.1:{port}", str(tpl)
    srv.should_exit = True
    thread.join(timeout=5)


def _json(result):
    for block in result.content:
        text = getattr(block, "text", "")
        if text.strip().startswith("{"):
            return json.loads(text)
    structured = getattr(result, "structuredContent", None)
    if structured:
        return structured.get("result", structured)   # FastMCP nests dict returns under "result"
    raise AssertionError(f"no json in tool result: {result}")


QUESTIONS = [{"question": "Who is the audience?", "options": ["Classmates", "Faculty"]}]


async def _answer_questions(session):
    presented = await session.call_tool("present_deck_questions", {"questions": QUESTIONS})
    set_id = next(
        str(block.resource.uri).rsplit("/", 1)[-1]
        for block in presented.content if getattr(block, "type", "") == "resource"
    )
    submitted = await session.call_tool("submit_deck_answers", {
        "set_id": set_id,
        "answers": [{"question": "Who is the audience?", "answer": "Classmates"}],
    })
    assert "error" not in _json(submitted)


async def _make_deck(base, user, template):
    from mcp.client.streamable_http import streamablehttp_client
    from mcp.client.session import ClientSession
    headers = {"X-Gies-Key": "testkey", "X-Gies-User": user}
    async with streamablehttp_client(f"{base}/mcp", headers=headers) as (r, w, _):
        async with ClientSession(r, w) as session:
            await session.initialize()
            blocked = await session.call_tool(
                "create_presentation_from_template", {"template_path": template})
            assert "questions" in _json(blocked)["error"]      # gate shut until answered
            await _answer_questions(session)
            created = await session.call_tool(
                "create_presentation_from_template", {"template_path": template})
            pid = _json(created)["presentation_id"]
            await session.call_tool(
                "add_slide", {"layout_index": 0, "presentation_id": pid})
            saved = await session.call_tool(
                "save_presentation", {"file_path": f"{user}.pptx", "presentation_id": pid})
            return _json(saved)["download_url"]


@pytest.mark.asyncio
async def test_round_trip_downloads_valid_pptx(server, tmp_path):
    base, template = server
    url = await _make_deck(base, "alice", template)
    async with httpx.AsyncClient() as c:
        r = await c.get(url)
    assert r.status_code == 200
    assert "presentationml" in r.headers["content-type"]
    out = tmp_path / "got.pptx"
    out.write_bytes(r.content)
    assert len(Presentation(str(out)).slides) == 1


@pytest.mark.asyncio
async def test_users_are_isolated(server, tmp_path):
    """Same guessable id 'presentation_1' must resolve to each user's own deck."""
    base, template = server
    # Alice builds from the template (its slide layouts) and adds one slide.
    alice_url = await _make_deck(base, "alice", template)
    # Bob builds a blank deck with no template — his presentation_1 is different.
    from mcp.client.streamable_http import streamablehttp_client
    from mcp.client.session import ClientSession
    async with streamablehttp_client(
        f"{base}/mcp", headers={"X-Gies-Key": "testkey", "X-Gies-User": "bob"}
    ) as (r, w, _):
        async with ClientSession(r, w) as session:
            await session.initialize()
            await _answer_questions(session)
            created = await session.call_tool("create_presentation", {})
            pid = _json(created)["presentation_id"]
            assert pid == "presentation_1"       # same id string as alice's first deck
            saved = await session.call_tool(
                "save_presentation", {"file_path": "bob.pptx", "presentation_id": pid})
            bob_url = _json(saved)["download_url"]
    async with httpx.AsyncClient() as c:
        bob = await c.get(bob_url)
        alice = await c.get(alice_url)
    bob_deck = tmp_path / "bob.pptx"; bob_deck.write_bytes(bob.content)
    alice_deck = tmp_path / "alice.pptx"; alice_deck.write_bytes(alice.content)
    assert len(Presentation(str(bob_deck)).slides) == 0     # bob's blank deck
    assert len(Presentation(str(alice_deck)).slides) == 1   # alice's, untouched


@pytest.mark.asyncio
async def test_uploaded_design_builds_deck(server, tmp_path):
    """Full custom-template path: upload card → HTTP POST → upload_ready →
    create from the uploaded file → download → reparse."""
    import io
    from pptx import Presentation as P
    from mcp.client.streamable_http import streamablehttp_client
    from mcp.client.session import ClientSession
    base, _ = server
    headers = {"X-Gies-Key": "testkey", "X-Gies-User": "carol"}
    design = io.BytesIO()
    P().save(design)
    async with streamablehttp_client(f"{base}/mcp", headers=headers) as (r, w, _sid):
        async with ClientSession(r, w) as session:
            await session.initialize()
            presented = await session.call_tool("present_upload_card", {})
            upload_url = next(
                str(block.resource.uri) for block in presented.content
                if getattr(block, "type", "") == "resource"
            ).replace("ui://pptx", base)
            async with httpx.AsyncClient() as c:
                posted = await c.post(
                    f"{upload_url}?name=my-course-design.pptx", content=design.getvalue())
            assert posted.status_code == 200
            token = upload_url.rsplit("/", 1)[-1]
            info = await session.call_tool("upload_ready", {"upload_id": token})
            file_name = _json(info)["file_name"]
            assert file_name == "upload-my-course-design.pptx"
            await _answer_questions(session)
            created = await session.call_tool(
                "create_presentation_from_template", {"template_path": file_name})
            pid = _json(created)["presentation_id"]
            saved = await session.call_tool(
                "save_presentation", {"file_path": "from-design.pptx", "presentation_id": pid})
            url = _json(saved)["download_url"]
    async with httpx.AsyncClient() as c:
        got = await c.get(url)
    assert got.status_code == 200
    out = tmp_path / "from-design.pptx"
    out.write_bytes(got.content)
    P(str(out))                                       # parses — built from the upload


@pytest.mark.asyncio
async def test_missing_auth_rejected(server):
    base, _ = server
    async with httpx.AsyncClient() as c:
        r = await c.post(f"{base}/mcp", json={"jsonrpc": "2.0", "id": 1, "method": "initialize"})
    assert r.status_code == 401
