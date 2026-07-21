import io

import httpx
import pytest
from pptx import Presentation
from starlette.applications import Starlette
from starlette.routing import Route

import gies_auth
import gies_sandbox as sb
import gies_uploads as up


def _pptx_bytes() -> bytes:
    """A deck WITH content slides — uploads must be stripped to a design shell."""
    pres = Presentation()
    for title in ("Old project intro", "Old project data"):
        slide = pres.slides.add_slide(pres.slide_layouts[0])
        slide.placeholders[0].text_frame.text = title
    buf = io.BytesIO()
    pres.save(buf)
    return buf.getvalue()


@pytest.fixture(autouse=True)
def _ctx(tmp_path, monkeypatch):
    monkeypatch.setattr(sb, "SANDBOX_ROOT", tmp_path / "decks")
    up._tokens.clear()
    up._completed.clear()
    gies_auth._user.set("alice")


def _client():
    app = Starlette(routes=[Route("/upload/{token}", up.upload, methods=["POST", "OPTIONS"])])
    return httpx.AsyncClient(transport=httpx.ASGITransport(app), base_url="http://t")


async def test_valid_upload_lands_in_sandbox_and_consumes_token():
    token = up.mint("alice")
    async with _client() as c:
        r = await c.post(f"/upload/{token}?name=My Design (v2).pptx", content=_pptx_bytes())
    assert r.status_code == 200
    payload = r.json()
    name = payload["file_name"]
    assert name.startswith("upload-") and name.endswith(".pptx")
    assert payload["slides_removed"] == 2                 # old content stripped
    saved = sb.user_root("alice") / name
    assert saved.exists()
    shell = Presentation(str(saved))
    assert len(shell.slides) == 0                         # design shell: no slides
    assert len(shell.slide_layouts) > 0                   # ...but layouts survive
    async with _client() as c:                       # token is one-shot
        r2 = await c.post(f"/upload/{token}?name=x.pptx", content=_pptx_bytes())
    assert r2.status_code == 404


async def test_bad_payloads_keep_token_alive():
    token = up.mint("alice")
    async with _client() as c:
        r = await c.post(f"/upload/{token}?name=notes.pptx", content=b"not a pptx at all")
        assert r.status_code == 400
        big = b"x" * (up.MAX_BYTES + 1)
        r2 = await c.post(f"/upload/{token}?name=big.pptx", content=big)
        assert r2.status_code == 413
        r3 = await c.post(f"/upload/{token}?name=ok.pptx", content=_pptx_bytes())
    assert r3.status_code == 200                     # token survived both failures


async def test_expired_and_unknown_tokens_rejected(monkeypatch):
    async with _client() as c:
        assert (await c.post("/upload/nope?name=a.pptx", content=_pptx_bytes())).status_code == 404
    monkeypatch.setattr(up, "TOKEN_TTL_SECONDS", -1)
    token = up.mint("alice")
    async with _client() as c:
        assert (await c.post(f"/upload/{token}?name=a.pptx", content=_pptx_bytes())).status_code == 404


async def test_preflight_cors():
    async with _client() as c:
        r = await c.options("/upload/whatever")
    assert r.status_code == 204
    assert r.headers["access-control-allow-origin"] == "*"
    assert "POST" in r.headers["access-control-allow-methods"]


async def test_upload_ready_scoped_to_user():
    token = up.mint("alice")
    async with _client() as c:
        await c.post(f"/upload/{token}?name=design.pptx", content=_pptx_bytes())
    gies_auth._user.set("bob")
    assert "error" in up.ready(token)
    gies_auth._user.set("alice")
    info = up.ready(token)
    assert "error" not in info
    assert info["file_name"].startswith("upload-")
    assert info["slide_count"] == 0
    assert isinstance(info["layouts"], list) and len(info["layouts"]) > 0
    assert {"index", "name", "placeholders"} <= set(info["layouts"][0])


def test_ready_without_upload_errors():
    assert "error" in up.ready("never-uploaded")


def test_card_contains_endpoint_and_action():
    token = up.mint("alice")
    card = up.render_card(token)
    assert f"/upload/{token}" in card
    assert "upload_ready" in card
