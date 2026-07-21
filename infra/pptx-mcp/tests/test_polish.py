import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

import gies_auth
import gies_sandbox as sb
from utils import presentation_utils as pu


@pytest.fixture(autouse=True)
def _ctx(tmp_path, monkeypatch):
    monkeypatch.setattr(sb, "SANDBOX_ROOT", tmp_path / "decks")
    gies_auth._user.set("alice")


def _build_messy_deck():
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[1])   # title + content, both left empty
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "tiny body text"
    run.font.size = Pt(11)
    return pres


def test_save_strips_empty_placeholders_and_bumps_fonts():
    saved = pu.save_presentation(_build_messy_deck(), "deck.pptx")
    slide = Presentation(saved).slides[0]
    assert len(slide.placeholders) == 0                    # no "Click to add ..." prompts
    run = slide.shapes[0].text_frame.paragraphs[0].runs[0]
    assert run.text == "tiny body text"
    assert run.font.size >= Pt(pu.MIN_FONT_PT)


def test_save_keeps_filled_placeholders():
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[1])
    slide.placeholders[0].text_frame.text = "Real Title"
    saved = pu.save_presentation(pres, "titled.pptx")
    out = Presentation(saved).slides[0]
    kept = [p.text_frame.text for p in out.placeholders]
    assert kept == ["Real Title"]                          # filled stays, empty body removed
