"""Decks attached in the chat composer become the deck's design.

LibreChat's backend POSTs the bytes here directly - server to server, because a
real template runs to megabytes and could never travel as a tool argument. The
route keeps header auth (unlike token-authenticated /upload) since the caller
holds the shared secret.
"""

import io

import httpx
import pytest
from pptx import Presentation
from starlette.applications import Starlette
from starlette.routing import Route

import gies_auth
import gies_sandbox as sb
import gies_uploads as up


def _pptx_bytes(slides=("Old content one", "Old content two")) -> bytes:
    pres = Presentation()
    for title in slides:
        slide = pres.slides.add_slide(pres.slide_layouts[0])
        slide.placeholders[0].text_frame.text = title
    buf = io.BytesIO()
    pres.save(buf)
    return buf.getvalue()


@pytest.fixture(autouse=True)
def _ctx(tmp_path, monkeypatch):
    monkeypatch.setattr(sb, "SANDBOX_ROOT", tmp_path / "decks")
    up._tokens.clear()
    up._completed.clear()
    up._attached.clear()
    gies_auth._user.set("alice")


def _client():
    app = Starlette(routes=[Route("/design", up.design_upload, methods=["POST"])])
    return httpx.AsyncClient(transport=httpx.ASGITransport(app), base_url="http://t")


async def test_attached_deck_is_stripped_and_registered():
    async with _client() as c:
        r = await c.post("/design?name=course-template.pptx", content=_pptx_bytes())
    assert r.status_code == 200
    body = r.json()
    assert body["slides_removed"] == 2
    assert body["file_name"] == "upload-course-template.pptx"
    assert body["design_layouts"][0]["slides_used"] == 2
    assert "alice" in up._attached


async def test_original_slides_never_leak():
    async with _client() as c:
        await c.post("/design?name=t.pptx", content=_pptx_bytes())
    saved = Presentation(up._attached["alice"][0])
    assert len(saved.slides) == 0


async def test_tool_returns_the_registered_design():
    async with _client() as c:
        await c.post("/design?name=t.pptx", content=_pptx_bytes())
    result = up.attached("alice")
    assert result["file_name"] == "upload-t.pptx"
    assert result["slide_count"] == 0
    assert result["design_layouts"]
    assert "create_presentation_from_template" in result["message"]


async def test_other_users_cannot_see_it():
    async with _client() as c:
        await c.post("/design?name=t.pptx", content=_pptx_bytes())
    assert "error" in up.attached("bob")


async def test_nothing_attached_points_at_the_upload_card():
    result = up.attached("alice")
    assert "present_upload_card" in result["error"]


async def test_non_pptx_is_rejected():
    async with _client() as c:
        r = await c.post("/design?name=t.pptx", content=b"definitely not a deck")
    assert r.status_code == 400
    assert "alice" not in up._attached


async def test_oversize_is_rejected_before_reading():
    async with _client() as c:
        r = await c.post(
            "/design?name=t.pptx",
            content=b"x",
            headers={"content-length": str(up.MAX_BYTES + 1)},
        )
    assert r.status_code == 413


async def test_reattaching_replaces_the_previous_design():
    async with _client() as c:
        await c.post("/design?name=first.pptx", content=_pptx_bytes(("a",)))
        await c.post("/design?name=second.pptx", content=_pptx_bytes(("b", "c", "d")))
    assert up.attached("alice")["file_name"] == "upload-second.pptx"
