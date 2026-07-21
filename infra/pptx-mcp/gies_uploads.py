"""User-supplied design/template upload for the deck builder.

The model presents an upload card (ui:// resource, rendered in the dock); the
browser POSTs the .pptx straight to /upload/<token> on this container — the
LibreChat file store is never involved. The token is the credential (the route
is exempt from header auth, mirroring /download in reverse): one-shot, per-user,
short TTL, consumed only on a successful save so failed attempts can retry.
Files land in the uploader's sandbox dir, where the existing
create_presentation_from_template picks them up with no build-flow changes.
"""
import io
import re
import time
from typing import Dict, Tuple

from pptx import Presentation
from pptx.oxml.ns import qn
from starlette.requests import Request
from starlette.responses import JSONResponse, PlainTextResponse, Response
from mcp.server.fastmcp import FastMCP
from mcp.types import EmbeddedResource, TextResourceContents

import gies_sandbox
from gies_auth import current_user
from gies_downloads import PUBLIC_URL

TOKEN_TTL_SECONDS = 15 * 60
MAX_BYTES = 15 * 1024 * 1024

_tokens: Dict[str, Tuple[str, float]] = {}      # token -> (user, expires_at)
_completed: Dict[str, Tuple[str, str]] = {}     # token -> (user, saved_path)

_CORS = {
    "Access-Control-Allow-Origin": "*",
    "Access-Control-Allow-Methods": "POST, OPTIONS",
    "Access-Control-Allow-Headers": "content-type",
}


def mint(user: str) -> str:
    import secrets

    token = secrets.token_urlsafe(24)
    _tokens[token] = (user, time.time() + TOKEN_TTL_SECONDS)
    return token


def _strip_slides(pres: Presentation) -> int:
    """Reduce an uploaded deck to a design shell: drop every slide, keep the
    masters, layouts, theme, and fonts. 'Use my deck as the design' means its
    look — never its old content leaking into the generated deck."""
    slide_ids = pres.slides._sldIdLst
    removed = 0
    for slide_id in list(slide_ids):
        pres.part.drop_rel(slide_id.get(qn("r:id")))
        slide_ids.remove(slide_id)
        removed += 1
    return removed


def _safe_name(raw: str) -> str:
    stem = re.sub(r"[^\w.-]+", "-", raw.rsplit("/", 1)[-1]).strip("-.") or "design"
    if stem.lower().endswith(".pptx"):
        stem = stem[:-5]
    return f"upload-{stem[:60]}.pptx"


async def upload(request: Request) -> Response:
    if request.method == "OPTIONS":
        return Response(status_code=204, headers=_CORS)

    token = request.path_params["token"]
    entry = _tokens.get(token)
    if entry is None or entry[1] < time.time():
        _tokens.pop(token, None)
        return PlainTextResponse(
            "This upload link has expired — ask GiesChat for a fresh upload card.",
            status_code=404, headers=_CORS,
        )
    user, _ = entry

    too_large = PlainTextResponse(
        f"File is too large — the limit is {MAX_BYTES // (1024 * 1024)} MB.",
        status_code=413, headers=_CORS,
    )
    declared = request.headers.get("content-length")
    if declared is not None and declared.isdigit() and int(declared) > MAX_BYTES:
        return too_large
    chunks: list = []
    received = 0
    async for chunk in request.stream():
        received += len(chunk)
        if received > MAX_BYTES:
            return too_large
        chunks.append(chunk)
    body = b"".join(chunks)
    try:
        parsed = Presentation(io.BytesIO(body))
    except Exception:
        return PlainTextResponse(
            "That file could not be read as a PowerPoint (.pptx) presentation.",
            status_code=400, headers=_CORS,
        )

    file_name = _safe_name(request.query_params.get("name", "design.pptx"))
    path = gies_sandbox.resolve(file_name, user)
    slides_removed = _strip_slides(parsed)
    parsed.save(path)

    _tokens.pop(token, None)
    _completed[token] = (user, path)
    return JSONResponse(
        {"file_name": file_name, "slides_removed": slides_removed,
         "layouts": len(parsed.slide_layouts)},
        headers=_CORS,
    )


def ready(upload_id: str) -> Dict:
    entry = _completed.get(upload_id)
    if entry is None or entry[0] != current_user():
        return {"error": "No completed upload found for this id. Re-present the upload card with present_upload_card."}
    _, path = entry
    try:
        pres = Presentation(path)
    except Exception:
        return {"error": "The uploaded file is no longer readable. Ask the user to upload it again."}
    file_name = path.rsplit("/", 1)[-1]
    return {
        "file_name": file_name,
        "slide_count": len(pres.slides),
        "layouts": [
            {"index": i, "name": layout.name, "placeholders": len(layout.placeholders)}
            for i, layout in enumerate(pres.slide_layouts)
        ],
        "message": (
            f"Design uploaded and reduced to an empty shell of its layouts — its "
            f"original slides were removed. Build EVERY slide of the new deck "
            f"yourself with create_presentation_from_template using template_path "
            f"\"{file_name}\" and the layouts listed above."
        ),
    }


def render_card(token: str) -> str:
    endpoint = f"{PUBLIC_URL}/upload/{token}"
    return (
        CARD_TEMPLATE
        .replace("__ENDPOINT__", endpoint)
        .replace("__TOKEN__", token)
        .replace("__MAX_MB__", str(MAX_BYTES // (1024 * 1024)))
    )


def register_upload_tools(app: FastMCP) -> None:
    @app.tool()
    def present_upload_card() -> list:
        """Show the user a card to upload their own .pptx design/outline. Wait for
        the card to trigger upload_ready — never call upload_ready yourself."""
        token = mint(current_user())
        return [
            EmbeddedResource(
                type="resource",
                resource=TextResourceContents(
                    uri=f"ui://pptx/upload/{token}",
                    mimeType="text/html",
                    text=render_card(token),
                ),
            ),
            {
                "type": "text",
                "text": "Upload card presented. Stop and wait for the user's upload; do not build yet.",
            },
        ]

    @app.tool()
    def upload_ready(upload_id: str) -> Dict:
        """Called by the upload card once the file is stored. Returns the uploaded
        design's file name and layout inventory for building."""
        return ready(upload_id)


CARD_TEMPLATE = """<!doctype html><html><head><meta charset="utf-8"><style>
  body { margin: 0; background: transparent;
         font: 14.5px/1.5 system-ui, -apple-system, "Segoe UI", sans-serif; }
  .card { background: #13294B; color: #edf1f7; border-radius: 16px;
          padding: 12px 16px 12px;
          box-shadow: 0 6px 24px rgba(0,0,0,.35); }
  .title { font-family: "Iowan Old Style", Georgia, serif; font-size: 17.5px;
           font-weight: 400; margin: 0 0 10px; padding: 0 4px 0; }
  .zone { border: 1.5px dashed #2d4b7d; border-radius: 11px; padding: 17px 16px;
          text-align: center; color: #9fb1cc; cursor: pointer; font-size: 15px; }
  .zone:hover, .zone.drag { background: #1e3a66; color: #edf1f7; }
  .zone b { color: #edf1f7; font-weight: 600; }
  .hint { font-size: 13.5px; color: #9fb1cc; margin-top: 10px; padding: 0 6px; }
  .err { color: #e8927c; font-size: 14.5px; margin-top: 10px; padding: 0 6px; display: none; }
  .done { font-size: 15px; color: #9fb1cc; padding: 8px 6px; display: none; }
  .done b { color: #edf1f7; font-weight: 600; }
  input[type=file] { display: none; }
  .zone:focus-visible { outline: 2px solid #E84A27; outline-offset: 2px; }
</style></head><body>
<div class="card">
  <h2 class="title">Upload your slide design</h2>
  <div class="zone" id="zone" role="button" tabindex="0"
       aria-label="Upload a .pptx design">
    <b>Drop a .pptx here</b> or click to choose a file
  </div>
  <div class="hint">Your deck will be built using this file's layouts, colors, and fonts. Max __MAX_MB__ MB.</div>
  <div class="err" id="err"></div>
  <div class="done" id="done"><b>Design received</b> — building with your layouts.</div>
  <input type="file" id="file" accept=".pptx">
</div>
<script>
  var zone = document.getElementById("zone");
  var input = document.getElementById("file");
  var err = document.getElementById("err");

  function fail(message) { err.textContent = message; err.style.display = "block"; }

  function send(file) {
    err.style.display = "none";
    if (!file) { return; }
    if (!/\\.pptx$/i.test(file.name)) { fail("Please choose a .pptx file."); return; }
    zone.textContent = "Uploading " + file.name + "\\u2026";
    fetch("__ENDPOINT__?name=" + encodeURIComponent(file.name), { method: "POST", body: file })
      .then(function (r) {
        if (!r.ok) { return r.text().then(function (t) { throw new Error(t); }); }
        return r.json();
      })
      .then(function () {
        zone.style.display = "none";
        document.querySelector(".hint").style.display = "none";
        document.getElementById("done").style.display = "block";
        window.parent.postMessage({
          type: "tool",
          payload: { toolName: "upload_ready", params: { upload_id: "__TOKEN__" } }
        }, "*");
      })
      .catch(function (e) {
        zone.innerHTML = "<b>Drop a .pptx here</b> or click to choose a file";
        fail(e.message || "Upload failed — please try again.");
      });
  }

  zone.onclick = function () { input.click(); };
  zone.onkeydown = function (e) { if (e.key === "Enter" || e.key === " ") input.click(); };
  input.onchange = function () { send(input.files[0]); };
  zone.ondragover = function (e) { e.preventDefault(); zone.classList.add("drag"); };
  zone.ondragleave = function () { zone.classList.remove("drag"); };
  zone.ondrop = function (e) {
    e.preventDefault(); zone.classList.remove("drag");
    send(e.dataTransfer.files[0]);
  };
</script>
</body></html>"""
